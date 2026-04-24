"""Upwise angel pitch deck builder — full redesign.

Principles:
- One idea per slide, strong hierarchy.
- Hero headlines (28pt) sized to fit two-line max within fixed containers.
- matplotlib charts where they earn their slot (Slides 3 and 7).
- Coral used sparingly as an accent, never as background filler.
- Generous whitespace, grid-aligned layouts.
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# ───────────────────────── PALETTE ─────────────────────────
NAVY         = RGBColor(0x0F, 0x17, 0x2A)
OFFWHITE     = RGBColor(0xF8, 0xF9, 0xFA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CORAL        = RGBColor(0xF9, 0x61, 0x67)
LIGHT_CORAL  = RGBColor(0xFC, 0xA5, 0xAA)
CORAL_TINT   = RGBColor(0xFF, 0xEE, 0xEE)
SLATE        = RGBColor(0x47, 0x55, 0x6D)
MID_SLATE    = RGBColor(0x64, 0x74, 0x8B)
LIGHT_SLATE  = RGBColor(0x94, 0xA3, 0xB8)
PALE_SLATE   = RGBColor(0xCB, 0xD5, 0xE1)
LIGHT_GREY   = RGBColor(0xE5, 0xE7, 0xEB)

CORAL_HEX    = "#F96167"
SLATE_HEX    = "#47556D"
MID_SLATE_HEX = "#64748B"
LIGHT_SLATE_HEX = "#94A3B8"
OFFWHITE_HEX = "#F8F9FA"

# ───────────────────────── FONTS ─────────────────────────
FONT = "Inter"  # falls back to Arial/Calibri on host systems without Inter

# ───────────────────────── LAYOUT ─────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.7
MARGIN_R = 0.7
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # 11.933"

# ───────────────────────── PATHS ─────────────────────────
ROOT = "/sessions/charming-funny-hypatia/mnt/school/deck"
CHARTS = f"{ROOT}/charts"
os.makedirs(CHARTS, exist_ok=True)

# ───────────────────────── HELPERS ─────────────────────────

def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_line(shape):
    shape.line.fill.background()


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, *, fill=None, line_color=None, line_width=0.75, line_dash=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill is not None:
        set_fill(shape, fill)
    else:
        shape.fill.background()
    if line_color is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
        if line_dash is not None:
            shape.line.dash_style = line_dash
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, italic=False, color=SLATE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, font_name=FONT):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, side, Inches(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multi_text(slide, left, top, width, height, runs, *,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   line_spacing=1.2):
    """runs = list of dicts: {text, size, bold, italic, color, newline (bool)}"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, side, Inches(0))
    tf.vertical_anchor = anchor

    first = True
    current_p = tf.paragraphs[0]
    current_p.alignment = align
    current_p.line_spacing = line_spacing

    for run in runs:
        if run.get("newline") and not first:
            current_p = tf.add_paragraph()
            current_p.alignment = run.get("align", align)
            current_p.line_spacing = run.get("line_spacing", line_spacing)
            if run.get("space_before"):
                current_p.space_before = Pt(run["space_before"])
        r = current_p.add_run()
        r.text = run["text"]
        r.font.name = FONT
        r.font.size = Pt(run.get("size", 14))
        r.font.bold = run.get("bold", False)
        r.font.italic = run.get("italic", False)
        r.font.color.rgb = run.get("color", SLATE)
        first = False
    return tb


def add_circle(slide, cx, cy, diameter, fill, *, text=None, text_color=WHITE, text_size=22, bold=True):
    left = cx - diameter / 2
    top = cy - diameter / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter)
    )
    set_fill(shape, fill)
    no_line(shape)
    if text:
        tf = shape.text_frame
        for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, side, Inches(0))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.bold = bold
        r.font.size = Pt(text_size)
        r.font.color.rgb = text_color
    return shape


def add_line(slide, x1, y1, x2, y2, color=SLATE, width=1.0):
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def add_headline(slide, text, *, size=28, top=0.55, height=1.3, color=NAVY, width=None):
    """Standardised headline block. Fixed-size container; word wraps to 2 lines max."""
    if width is None:
        width = CONTENT_W
    return add_text(
        slide, MARGIN_L, top, width, height, text,
        size=size, bold=True, color=color, line_spacing=1.1
    )


def add_subhead(slide, text, *, top=1.65, size=17, color=MID_SLATE, italic=False, width=None):
    if width is None:
        width = CONTENT_W
    return add_text(
        slide, MARGIN_L, top, width, 0.5, text,
        size=size, italic=italic, color=color
    )


# ───────────────────────── MATPLOTLIB CHART RENDERERS ─────────────────────────

def style_axes(ax):
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color(LIGHT_SLATE_HEX)
        ax.spines[spine].set_linewidth(0.8)
    ax.tick_params(colors=SLATE_HEX, labelsize=12)


def render_market_chart(path):
    rcParams["font.family"] = "sans-serif"
    rcParams["font.sans-serif"] = ["Inter", "Arial", "Helvetica", "DejaVu Sans"]

    labels = ["Kumon (2 subjects)", "Cluey (weekly tutoring)", "Upwise family (up to 4 kids)",
              "Mathletics", "Reading Eggs"]
    values = [3840, 2500, 708, 170, 110]
    colors = [MID_SLATE_HEX, MID_SLATE_HEX, CORAL_HEX, MID_SLATE_HEX, MID_SLATE_HEX]

    fig, ax = plt.subplots(figsize=(11.5, 3.6), dpi=180)
    fig.patch.set_facecolor(OFFWHITE_HEX)
    ax.set_facecolor(OFFWHITE_HEX)

    bars = ax.barh(labels[::-1], values[::-1], color=colors[::-1], height=0.6, zorder=3)
    ax.invert_yaxis()
    ax.grid(axis="x", color="#DDE0E6", linewidth=0.6, zorder=0)

    for bar, v in zip(bars, values[::-1]):
        ax.text(
            v + max(values) * 0.015, bar.get_y() + bar.get_height() / 2,
            f"A${v:,}",
            va="center", ha="left",
            fontsize=14, fontweight="bold",
            color=CORAL_HEX if v == 708 else SLATE_HEX
        )

    ax.set_xlabel("Annual cost per family (A$)", fontsize=12, color=SLATE_HEX, labelpad=8)
    ax.tick_params(axis="y", labelsize=13, colors=SLATE_HEX)
    ax.tick_params(axis="x", labelsize=11, colors=MID_SLATE_HEX)
    style_axes(ax)
    ax.set_xlim(0, 4400)

    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor=OFFWHITE_HEX)
    plt.close(fig)


def render_evidence_chart(path):
    rcParams["font.family"] = "sans-serif"
    rcParams["font.sans-serif"] = ["Inter", "Arial", "Helvetica", "DejaVu Sans"]

    labels = ["Bloom\n(1984)", "Kulik & Kulik\n(1990)", "RAND / Carnegie\n(2014)", "AI Tutoring Review\n(2025)"]
    values = [2.00, 0.52, 0.43, 0.65]
    colors = [CORAL_HEX, MID_SLATE_HEX, MID_SLATE_HEX, MID_SLATE_HEX]

    fig, ax = plt.subplots(figsize=(11.5, 4.2), dpi=180)
    fig.patch.set_facecolor(OFFWHITE_HEX)
    ax.set_facecolor(OFFWHITE_HEX)

    bars = ax.bar(labels, values, color=colors, width=0.55, zorder=3)
    ax.grid(axis="y", color="#DDE0E6", linewidth=0.6, zorder=0)

    for bar, v in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2, v + 0.06,
            f"+{v:.2f} SD",
            ha="center", va="bottom",
            fontsize=14, fontweight="bold",
            color=CORAL_HEX if v == 2.0 else SLATE_HEX
        )

    ax.set_ylabel("Effect size (standard deviations)", fontsize=12, color=SLATE_HEX, labelpad=10)
    ax.tick_params(axis="x", labelsize=12, colors=SLATE_HEX)
    ax.tick_params(axis="y", labelsize=11, colors=MID_SLATE_HEX)
    style_axes(ax)
    ax.set_ylim(0, 2.4)

    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor=OFFWHITE_HEX)
    plt.close(fig)


# ───────────────────────── MAIN ─────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    market_chart = f"{CHARTS}/market-chart.png"
    evidence_chart = f"{CHARTS}/evidence-chart.png"
    render_market_chart(market_chart)
    render_evidence_chart(evidence_chart)

    # ───── Slide 1 — Cover ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_text(s, 0.9, 2.9, 11.5, 1.4, "Upwise",
             size=84, bold=True, color=WHITE, line_spacing=1.0)
    add_rect(s, 0.95, 4.35, 1.3, 0.06, fill=CORAL)
    add_text(s, 0.9, 4.55, 11.5, 0.7,
             "Strong foundations in maths and English — built session by session.",
             size=22, color=PALE_SLATE)
    add_text(s, 0.9, 6.8, 11.5, 0.3,
             "Jesse, Founder · April 2026 · upwise.com.au",
             size=12, color=LIGHT_SLATE)
    add_notes(s, "Hi. I'm Jesse, founder of Upwise. In the next ten minutes I want to show you a business that I think is the clearest consumer opportunity in Australian education in a generation. Let's start with the problem every Australian parent already knows.")

    # ───── Slide 2 — The Problem ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "In a class of 28, your child doesn't stand a chance of being seen.",
                 size=26, top=0.55, height=1.5)

    card_top = 2.5
    card_h = 2.6
    gap = 0.3
    card_w = (CONTENT_W - 2 * gap) / 3

    stats = [
        ("42%", "of Australian schools teach with a qualified teacher shortage — up from 14% in 2018.", CORAL),
        ("82%", "of primary teachers say their job damages their mental health.", SLATE),
        ("1 in 3", "Year 3 students sit below literacy proficiency.", SLATE),
    ]
    for i, (number, label, num_color) in enumerate(stats):
        left = MARGIN_L + i * (card_w + gap)
        add_rect(s, left, card_top, card_w, card_h, fill=WHITE)
        add_text(s, left + 0.25, card_top + 0.25, card_w - 0.5, 1.2, number,
                 size=60, bold=True, color=num_color, align=PP_ALIGN.LEFT, line_spacing=1.0)
        add_text(s, left + 0.25, card_top + 1.55, card_w - 0.5, card_h - 1.6, label,
                 size=13, color=SLATE, line_spacing=1.35)

    band_top = 5.55
    band_h = 1.3
    add_rect(s, MARGIN_L, band_top, CONTENT_W, band_h, fill=CORAL_TINT)
    add_rect(s, MARGIN_L, band_top, 0.08, band_h, fill=CORAL)
    add_multi_text(
        s, MARGIN_L + 0.35, band_top + 0.2, CONTENT_W - 0.7, band_h - 0.3,
        [
            {"text": "Even Victoria's A$1.2 billion catch-up tutoring program ", "size": 14, "color": SLATE},
            {"text": "\"did not significantly improve outcomes\"", "size": 14, "color": SLATE, "italic": True},
            {"text": " — Victorian Auditor-General, 2024.", "size": 14, "color": SLATE},
            {"text": "The problem isn't money. It's the delivery model.", "size": 16, "bold": True, "color": NAVY, "newline": True, "space_before": 6},
        ]
    )
    add_notes(s, "Imagine a parent sitting at the kitchen table. Their Year 3 child hands them a report card that says 'Developing — needs additional support.' They call the school. The teacher is managing 28 kids with one aide and has no time for intervention. 35% of Australian parents already pay for tutoring to close the gap. They're not buying ambition — they're buying reassurance that their child won't be left behind. And when the government itself tried to fix it with a billion-dollar catch-up program, it didn't move the needle. The problem isn't money. It's the delivery model.")

    # ───── Slide 3 — The Market ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "A$2.2 billion is already leaving Australian parents' wallets.",
                 size=26, top=0.55, height=1.2)
    add_subhead(s, "Annual cost per family across Australian supplemental learning:",
                top=1.75, size=15, color=MID_SLATE)

    s.shapes.add_picture(
        f"{CHARTS}/market-chart.png",
        Inches(MARGIN_L), Inches(2.4),
        width=Inches(CONTENT_W)
    )

    add_text(s, MARGIN_L, 6.55, CONTENT_W, 0.4,
             "35% of Australian parents already pay for tutoring. 18% spend A$100–150 per week.",
             size=14, bold=True, color=SLATE, align=PP_ALIGN.LEFT)
    add_text(s, MARGIN_L, 6.95, CONTENT_W, 0.3,
             "Our job is redirection, not creation.",
             size=13, italic=True, color=MID_SLATE, align=PP_ALIGN.LEFT)
    add_notes(s, "Note what this slide is not. It's not a hockey-stick forecast or a TAM calculation from a McKinsey report. It's the money Australian parents are spending right now, today, on supplemental learning that isn't built for them. Kumon is nearly five times our family plan. Cluey is three times. Our job is redirection, not creation.")

    # ───── Slide 4 — Why Now ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "The 2-sigma solution, finally economical.",
                 size=30, top=0.55, height=1.0)
    add_subhead(s, "Three forces converging, for the first time in history.",
                top=1.55, size=16, italic=True)

    col_top = 2.7
    col_gap = 0.35
    col_w = (CONTENT_W - 2 * col_gap) / 3

    forces = [
        ("1", "The science is 40 years old.",
         "Benjamin Bloom (1984) proved 1:1 mastery tutoring produces +2.0 standard deviations of improvement — the 98th percentile of classroom-taught peers."),
        ("2", "The economics just flipped.",
         "LLM inference costs fell 150–1000× in five years. Personalised adaptive tutoring just became viable at a consumer subscription price."),
        ("3", "Parents are ready.",
         "Post-COVID, Australian parents normalised digital-first learning for their children — and are actively looking for tools that work."),
    ]
    for i, (num, title, desc) in enumerate(forces):
        left = MARGIN_L + i * (col_w + col_gap)
        circle_d = 0.8
        add_circle(s, left + circle_d / 2, col_top + circle_d / 2, circle_d,
                   CORAL, text=num, text_size=28, text_color=WHITE)
        add_text(s, left, col_top + 1.05, col_w, 0.7, title,
                 size=18, bold=True, color=NAVY, line_spacing=1.15)
        add_text(s, left, col_top + 1.85, col_w, 2.4, desc,
                 size=13, color=SLATE, line_spacing=1.4)

    add_text(s, MARGIN_L, 6.6, CONTENT_W, 0.5,
             "What Bloom called the \"2-sigma problem\" — the known solution that couldn't scale — is no longer a problem.",
             size=13, italic=True, color=MID_SLATE, line_spacing=1.3)
    add_notes(s, "Bloom's 2-Sigma Problem, 1984: one-on-one mastery tutoring produces outcomes classroom teaching can't match. We've known this for forty years. We couldn't do anything about it because one tutor per child wasn't economical. That changed, quietly, in the last eighteen months.")

    # ───── Slide 5 — What Upwise Does ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "A mastery-based learning coach for every Australian child.",
                 size=26, top=0.55, height=1.2)
    add_subhead(s, "For the price of a family streaming subscription.",
                top=1.8, size=17, italic=True)

    flow_top = 3.1
    circle_d = 1.0
    steps = [
        ("1", "Diagnose", "A short adaptive assessment pinpoints exactly where the gaps are."),
        ("2", "Learn", "20-minute daily sessions adapt in real time — questions stay in the 80% win-rate zone."),
        ("3", "Master", "Skills only lock in after proven mastery across multiple sessions. Spaced repetition holds them."),
    ]
    step_w = 3.3
    total_flow_w = 3 * step_w + 2 * 1.0
    flow_start = (SLIDE_W - total_flow_w) / 2

    for i, (num, title, desc) in enumerate(steps):
        left = flow_start + i * (step_w + 1.0)
        cx = left + step_w / 2
        add_circle(s, cx, flow_top + circle_d / 2, circle_d,
                   CORAL, text=num, text_size=36, text_color=WHITE)
        add_text(s, left, flow_top + circle_d + 0.2, step_w, 0.5, title,
                 size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.1)
        add_text(s, left, flow_top + circle_d + 0.8, step_w, 1.6, desc,
                 size=13, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.4)

    for i in range(2):
        x1 = flow_start + (i + 1) * step_w + i * 1.0
        x2 = x1 + 1.0
        cy = flow_top + circle_d / 2
        add_line(s, x1 + 0.05, cy, x2 - 0.05, cy, color=LIGHT_SLATE, width=1.8)

    strip_top = 6.4
    strip_h = 0.75
    add_rect(s, MARGIN_L, strip_top, CONTENT_W, strip_h, fill=NAVY)
    add_text(s, MARGIN_L, strip_top, CONTENT_W, strip_h,
             "Free 7-day trial  ·  A$39/mo single child  ·  A$59/mo family (up to 4 kids)  ·  ACARA-aligned",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(s, "Twenty minutes a day. Adaptive. Mastery-based. The kid can't jump ahead until they've actually understood. Which means every session builds on solid ground — no accumulating gaps. Free trial, then a subscription that costs less than a week of tutoring.")

    # ───── Slide 6 — The Parent Layer ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "The parent is the guide. We give them the map.",
                 size=30, top=0.55, height=1.0)
    add_subhead(s, "Research consistently shows AI tutoring works best when paired with a human guide.",
                top=1.6, size=15, italic=True)

    left_w = 6.8
    bullet_top = 2.7
    bullets = [
        ("Daily briefing", "Before each session — what your child's working on, what to watch for."),
        ("Real-time nudges", "When they get stuck — with the exact wording you can use."),
        ("Weekly reports", "Clear enough to share with their classroom teacher."),
        ("Conversation scripts", "What to say at dinner to reinforce what they just learned."),
    ]
    for i, (title, desc) in enumerate(bullets):
        row_top = bullet_top + i * 0.88
        add_rect(s, MARGIN_L, row_top + 0.15, 0.14, 0.14, fill=CORAL)
        add_multi_text(
            s, MARGIN_L + 0.35, row_top, left_w - 0.35, 0.8,
            [
                {"text": title, "size": 16, "bold": True, "color": NAVY},
                {"text": "  —  ", "size": 14, "color": LIGHT_SLATE},
                {"text": desc, "size": 14, "color": SLATE},
            ],
            line_spacing=1.35
        )

    quote_left = MARGIN_L + left_w + 0.4
    quote_w = CONTENT_W - left_w - 0.4
    quote_top = 2.7
    quote_h = 3.8
    add_rect(s, quote_left, quote_top, quote_w, quote_h, fill=NAVY)
    add_rect(s, quote_left, quote_top, quote_w, 0.12, fill=CORAL)
    add_text(s, quote_left + 0.4, quote_top + 0.6, quote_w - 0.8, 2.0,
             "\"No teaching degree required.\"",
             size=26, bold=True, color=WHITE, line_spacing=1.2)
    add_text(s, quote_left + 0.4, quote_top + 2.5, quote_w - 0.8, 1.2,
             "Every parent gets the tools. No extra hours. No curriculum expertise needed.",
             size=13, italic=True, color=PALE_SLATE, line_spacing=1.4)

    add_notes(s, "This is the deliberate bet. We're not replacing parents. We're not replacing teachers. We're giving parents the one thing they've always lacked — a map. What to say when their child is stuck. What to reinforce at dinner. Where the real gap is.")

    # ───── Slide 7 — The Evidence ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "This isn't experimental. It's proven. Four times over.",
                 size=26, top=0.55, height=1.0)

    s.shapes.add_picture(
        f"{CHARTS}/evidence-chart.png",
        Inches(MARGIN_L), Inches(1.9),
        width=Inches(CONTENT_W)
    )

    add_text(s, MARGIN_L, 6.3, CONTENT_W, 0.45,
             "Upwise is the first Australian operationalisation of 40 years of validated learning science.",
             size=15, bold=True, color=NAVY, line_spacing=1.3)
    add_text(s, MARGIN_L, 6.82, CONTENT_W, 0.35,
             "Sources: Bloom (1984) · Kulik & Kulik (1990) · Pane et al. / RAND (2014) · 2025 AI tutoring systematic review (48 studies).",
             size=11, italic=True, color=MID_SLATE, line_spacing=1.3)
    add_notes(s, "Our credibility isn't 'we built an AI.' It's a forty-year research lineage that four generations of learning scientists agree on. Bloom proved the 2-sigma effect in 1984. Kulik and Kulik meta-analysed it in 1990. RAND validated computer-based mastery in 2014. A 2025 systematic review of 48 AI tutoring studies confirms the strongest gains come when AI is paired with a human guide. Upwise is the first Australian operationalisation of that stack.")

    # ───── Slide 8 — Competition: The Empty Quadrant ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "No Australian product owns mastery + parent-as-guide.",
                 size=26, top=0.55, height=1.0)

    map_left = 1.8
    map_top = 2.1
    map_w = 9.7
    map_h = 4.2
    cx = map_left + map_w / 2
    cy = map_top + map_h / 2
    # Vertical axis — stops before the label area below
    add_line(s, cx, map_top - 0.1, cx, map_top + map_h, color=SLATE, width=1.5)
    add_line(s, map_left - 0.1, cy, map_left + map_w + 0.1, cy, color=SLATE, width=1.5)

    add_text(s, cx - 1.5, map_top - 0.45, 3.0, 0.3, "Parent-as-guide",
             size=12, bold=True, color=MID_SLATE, align=PP_ALIGN.CENTER)
    add_text(s, cx - 1.5, map_top + map_h + 0.18, 3.0, 0.3, "Parent-invisible",
             size=12, bold=True, color=MID_SLATE, align=PP_ALIGN.CENTER)
    add_text(s, map_left - 0.5, cy + 0.15, 2.5, 0.3, "Breadth-first",
             size=12, bold=True, color=MID_SLATE, align=PP_ALIGN.LEFT)
    add_text(s, map_left + map_w - 2.0, cy + 0.15, 2.5, 0.3, "Mastery-based",
             size=12, bold=True, color=MID_SLATE, align=PP_ALIGN.RIGHT)

    add_text(s, map_left + 0.4, map_top + 0.3, map_w / 2 - 0.5, 0.4,
             "(empty)", size=13, italic=True, color=LIGHT_SLATE)

    up_w = 3.0
    up_h = 1.0
    up_left = cx + (map_w / 2 - up_w) / 2
    up_top = map_top + 0.4
    add_rect(s, up_left, up_top, up_w, up_h, fill=CORAL)
    add_text(s, up_left, up_top, up_w, up_h, "UPWISE",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bl_names = ["Cluey", "Kumon", "Kip McGrath", "Matrix", "Kinetic Education"]
    for i, name in enumerate(bl_names):
        add_text(s, map_left + 0.4, cy + 0.35 + i * 0.35, 3.0, 0.3, name,
                 size=13, color=SLATE)

    br_names = ["Mathspace", "Reading Eggs", "Mathletics", "Prodigy", "Matific"]
    for i, name in enumerate(br_names):
        add_text(s, cx + 0.4, cy + 0.35 + i * 0.35, 3.0, 0.3, name,
                 size=13, color=SLATE)

    add_text(s, MARGIN_L, 6.75, CONTENT_W, 0.35,
             "Cluey FY25: revenue down 15% to A$25.6M, net loss A$5.5M  ·  Kumon: A$3,840/yr per child for maths + English",
             size=12, italic=True, color=MID_SLATE, align=PP_ALIGN.CENTER)
    add_notes(s, "Look at the quadrant Upwise sits in. Empty. Not because it's a bad quadrant — because no Australian incumbent has built for it. Cluey is losing money doing 1:1 human tutoring. Kumon is still running physical franchises in 2026. The IXLs and Mathletics of the world won't even look at you as a parent. That's our wedge.")

    # ───── Slide 9 — Business Model ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "17,000 families = A$10M ARR.",
                 size=34, top=0.55, height=0.9)
    add_subhead(s, "Not a moonshot. Arithmetic.",
                top=1.55, size=17, italic=True)

    hero_left = MARGIN_L
    hero_top = 2.5
    add_text(s, hero_left, hero_top, 5.0, 2.2, "1.1%",
             size=128, bold=True, color=CORAL, align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(s, hero_left, hero_top + 2.4, 5.2, 0.7,
             "of Australia's 1.52 million primary-aged households.",
             size=15, color=SLATE, line_spacing=1.35)
    add_text(s, hero_left, hero_top + 3.15, 5.2, 0.5,
             "Or 4.5% of the 380,000 already paying for supplemental learning.",
             size=13, italic=True, color=MID_SLATE, line_spacing=1.35)

    right_left = MARGIN_L + 5.5
    right_w = CONTENT_W - 5.5
    add_text(s, right_left, 2.5, right_w, 0.35, "PATH TO A$10M ARR",
             size=11, bold=True, color=CORAL)
    rows = [
        ("~17,100 active families", "at A$39–59/mo blend"),
        ("Software gross margin", "75–80% (LLM cost falling)"),
        ("Avg customer lifetime", "~18 months (assumption)"),
    ]
    for i, (a, b) in enumerate(rows):
        row_top = 2.95 + i * 0.55
        add_text(s, right_left, row_top, right_w * 0.55, 0.4, a,
                 size=14, bold=True, color=NAVY)
        add_text(s, right_left + right_w * 0.55, row_top, right_w * 0.45, 0.4, b,
                 size=13, color=MID_SLATE)

    add_text(s, right_left, 4.9, right_w, 0.35, "AU COMPARABLES",
             size=11, bold=True, color=CORAL)
    comps = [
        ("3P Learning (Mathletics):", "A$110M rev, software margins ✓", SLATE),
        ("Mathspace:", "profitable, bootstrapped ✓", SLATE),
        ("Cluey (1:1 human):", "A$25.6M rev, A$5.5M loss ✗", CORAL),
    ]
    for i, (a, b, color) in enumerate(comps):
        row_top = 5.35 + i * 0.42
        add_text(s, right_left, row_top, right_w * 0.45, 0.35, a,
                 size=12, bold=True, color=NAVY)
        add_text(s, right_left + right_w * 0.45, row_top, right_w * 0.55, 0.35, b,
                 size=12, color=color)

    add_text(s, MARGIN_L, 6.85, CONTENT_W, 0.4,
             "Cluey's loss isn't evidence AU edtech is broken. It's evidence that 1:1 human tutoring is.",
             size=13, italic=True, color=MID_SLATE, align=PP_ALIGN.CENTER)
    add_notes(s, "17,000 families. In a country with 1.52 million primary-age households, of which 380,000 already pay for tutoring. This is not a moonshot. This is arithmetic. And the unit economics work because we're software, not labour. Cluey's story is the proof — their A$5.5M loss isn't evidence AU edtech is broken. It's evidence that 1:1 human tutoring is.")

    # ───── Slide 10 — Safety & Trust ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "Safe AI. Australian-built. Parent-visible.",
                 size=30, top=0.55, height=1.0)
    add_subhead(s, "Every parent's first question about AI, answered in the product itself.",
                top=1.6, size=15, italic=True)

    cards_top = 2.8
    cards_h = 3.2
    gap = 0.15
    card_w = (CONTENT_W - 4 * gap) / 5

    cards = [
        ("Closed system", "No open chat. No user content. Curriculum-constrained."),
        ("20 min bounded", "Less screen time than YouTube. Designed to end on a win."),
        ("Parent-visible", "Every session observable. Nothing hidden."),
        ("Australian-owned", "AU data, AU servers. Techne AI Pty Ltd."),
        ("eSafety-aligned", "Designed to Australian eSafety Commissioner guidance."),
    ]
    for i, (title, desc) in enumerate(cards):
        left = MARGIN_L + i * (card_w + gap)
        add_rect(s, left, cards_top, card_w, cards_h, fill=WHITE)
        add_rect(s, left, cards_top, card_w, 0.12, fill=CORAL)
        add_text(s, left + 0.25, cards_top + 0.45, card_w - 0.5, 0.7, title,
                 size=15, bold=True, color=NAVY, line_spacing=1.2)
        add_text(s, left + 0.25, cards_top + 1.3, card_w - 0.5, cards_h - 1.4, desc,
                 size=12, color=SLATE, line_spacing=1.4)

    add_text(s, MARGIN_L, 6.5, CONTENT_W, 0.4,
             "No ads. No third-party tracking. No data sharing. Ever.",
             size=14, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_notes(s, "Every parent we speak to has the same question: is it safe? Is it just more screen time? Is it going to replace me as the parent? Our answer, on each: closed system, twenty minutes, parent-guided. Not a chatbot. Not an ad-driven app. Not a babysitter.")

    # ───── Slide 11 — Traction (placeholder) ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "Traction.", size=34, top=0.55, height=1.0)
    add_text(s, MARGIN_L, 1.7, CONTENT_W, 0.4,
             "Replace with actual numbers before pitching. Angels reward precision over puffery.",
             size=14, italic=True, color=MID_SLATE)

    grid_top = 2.5
    grid_h = 3.8
    grid_gap = 0.3
    cell_w = (CONTENT_W - grid_gap) / 2
    cell_h = (grid_h - grid_gap) / 2

    tiles = [
        ("Waitlist count",  "Growth rate per week"),
        ("Acquisition signal", "Organic or paid channel"),
        ("User engagement", "Retention, session count, NPS"),
        ("School / advisor interest", "LOIs, advisors, inbound press"),
    ]
    for i, (title, sub) in enumerate(tiles):
        row = i // 2
        col = i % 2
        left = MARGIN_L + col * (cell_w + grid_gap)
        top = grid_top + row * (cell_h + grid_gap)
        add_rect(s, left, top, cell_w, cell_h,
                 line_color=PALE_SLATE, line_width=1.5, line_dash=MSO_LINE_DASH_STYLE.DASH)
        add_text(s, left + 0.3, top + 0.4, cell_w - 0.6, 0.6, f"[{title}]",
                 size=20, bold=True, italic=True, color=LIGHT_SLATE)
        add_text(s, left + 0.3, top + 1.1, cell_w - 0.6, 0.5, sub,
                 size=13, italic=True, color=LIGHT_SLATE)

    add_notes(s, "Founder to draft based on actual numbers. Be honest — angels reward precision about early signals more than they reward puffed-up vanity metrics.")

    # ───── Slide 12 — Team (placeholder) ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, OFFWHITE)
    add_headline(s, "Team.", size=34, top=0.55, height=1.0)
    add_text(s, MARGIN_L, 1.7, CONTENT_W, 0.4,
             "Angels invest in people. Be specific about why you, why now, why this.",
             size=14, italic=True, color=MID_SLATE)

    col_top = 2.5
    col_h = 4.0
    col_gap = 0.4
    col_w = (CONTENT_W - 2 * col_gap) / 3

    cols = [
        ("Founder", "[Name + 1-line credential]", ""),
        ("Co-founder / Key hire", "[Name + credential]", "optional"),
        ("Advisors", "[Names + committed angels]", ""),
    ]
    for i, (role, name, note) in enumerate(cols):
        left = MARGIN_L + i * (col_w + col_gap)
        add_rect(s, left, col_top, col_w, col_h,
                 line_color=PALE_SLATE, line_width=1.5, line_dash=MSO_LINE_DASH_STYLE.DASH)
        add_rect(s, left + 0.4, col_top + 0.35, col_w - 0.8, 1.9,
                 fill=LIGHT_GREY)
        add_text(s, left + 0.4, col_top + 0.35, col_w - 0.8, 1.9,
                 "[Headshot]", size=14, italic=True, color=LIGHT_SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + 0.4, col_top + 2.45, col_w - 0.8, 0.4, f"[{role}]",
                 size=15, bold=True, italic=True, color=LIGHT_SLATE)
        add_text(s, left + 0.4, col_top + 2.9, col_w - 0.8, 0.5, name,
                 size=13, italic=True, color=LIGHT_SLATE)
        if note:
            add_text(s, left + 0.4, col_top + 3.4, col_w - 0.8, 0.3, note,
                     size=11, italic=True, color=LIGHT_SLATE)

    add_notes(s, "Founder-specific. Angels invest in people. Be specific about why you, why now, why this.")

    # ───── Slide 13 — The Ask ─────
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    left_x = 0.9
    left_w = 7.2
    add_text(s, left_x, 0.7, left_w, 1.0, "Raising A$[XXX]",
             size=44, bold=True, color=WHITE, line_spacing=1.1)
    add_text(s, left_x, 1.75, left_w, 0.5, "to reach A$[X]M ARR by [milestone]",
             size=18, italic=True, color=PALE_SLATE)

    add_text(s, left_x, 2.7, left_w, 0.35, "USE OF FUNDS",
             size=11, bold=True, color=CORAL)

    bar_top = 3.15
    bar_h = 0.55
    bar_w = left_w
    p_w = bar_w * 0.40
    g_w = bar_w * 0.40
    o_w = bar_w * 0.20
    add_rect(s, left_x, bar_top, p_w, bar_h, fill=CORAL)
    add_rect(s, left_x + p_w, bar_top, g_w, bar_h, fill=LIGHT_CORAL)
    add_rect(s, left_x + p_w + g_w, bar_top, o_w, bar_h, fill=MID_SLATE)

    add_text(s, left_x, bar_top + 0.12, p_w, 0.3, "Product  40%",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, left_x + p_w, bar_top + 0.12, g_w, 0.3, "Growth  40%",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, left_x + p_w + g_w, bar_top + 0.12, o_w, 0.3, "Ops  20%",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, left_x, 4.5, left_w, 0.35, "MILESTONES",
             size=11, bold=True, color=CORAL)
    milestones = [
        "→  Product launch",
        "→  Waitlist conversion (>60%)",
        "→  A$[X]k ARR",
        "→  CAC payback proof",
    ]
    for i, m in enumerate(milestones):
        add_text(s, left_x, 4.95 + i * 0.35, left_w, 0.32, m,
                 size=13, color=WHITE)

    add_text(s, left_x, 6.45, left_w, 0.3, "Runway: [X] months",
             size=13, italic=True, color=PALE_SLATE)

    right_x = 8.5
    right_w = SLIDE_W - right_x - 0.7
    add_text(s, right_x, 1.8, right_w, 3.5,
             "A generation of Australian children that never fall behind.",
             size=28, bold=True, color=CORAL, line_spacing=1.25)

    add_rect(s, right_x, 5.4, 1.0, 0.06, fill=CORAL)
    add_text(s, right_x, 5.55, right_w, 0.4,
             "That's what we're here to build.",
             size=14, italic=True, color=PALE_SLATE)

    add_text(s, 0, 7.1, SLIDE_W, 0.3,
             "jesse@activedigital.fund  ·  upwise.com.au",
             size=12, color=LIGHT_SLATE, align=PP_ALIGN.CENTER)

    add_notes(s, "Here's the ask. Here's what it funds. Here's what you get back — not in dollar terms, though I'm happy to walk through projections. This funds the build and the first twelve months of disciplined growth. A generation of Australian kids that don't fall behind. That's the thing worth building.")

    out = f"{ROOT}/upwise-angel-deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    build_presentation()
